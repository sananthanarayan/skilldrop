#!/usr/bin/env python3
"""Build a .pptx from a JSON deck spec.

The spec schema is documented in templates/deck-spec.json. Top-level fields:
- title:    overall deck title
- audience: archetype (informational; doesn't change rendering)
- palette:  either a hex-color block, or a string name resolvable from templates/palettes.json
- slides:   list of slide specs, each with a `layout` field

Supported layouts: title, section, content, two_column, big_number, quote, closing.

Usage:
    python3 build_deck.py <spec.json> [-o output.pptx]
"""

from __future__ import annotations

import argparse
import json
import sys
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

DEFAULT_PALETTE = {
    "primary": "#1A2A6C",
    "secondary": "#4A5A8C",
    "accent": "#FDBB2D",
    "background": "#FFFFFF",
    "text": "#222222",
}


def hex_color(s: str) -> RGBColor:
    return RGBColor.from_string(s.lstrip("#"))


def resolve_palette(p, script_dir: Path) -> dict:
    """Accepts a dict (use directly) or a string name (look up in templates/palettes.json)."""
    if isinstance(p, dict):
        return {**DEFAULT_PALETTE, **p}
    if isinstance(p, str):
        palettes_path = script_dir.parent / "templates" / "palettes.json"
        if not palettes_path.exists():
            sys.stderr.write(f"warning: palettes file missing at {palettes_path}; using defaults\n")
            return DEFAULT_PALETTE
        all_palettes = json.loads(palettes_path.read_text())
        if p not in all_palettes:
            sys.stderr.write(f"warning: palette {p!r} not found; using defaults\n")
            return DEFAULT_PALETTE
        return {**DEFAULT_PALETTE, **all_palettes[p]}
    return DEFAULT_PALETTE


def set_background(slide, hex_str: str) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = hex_color(hex_str)


def add_textbox(slide, text: str, *, left, top, width, height,
                size: int = 18, bold: bool = False, color: str = "#222222",
                align=PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = hex_color(color)


def add_bullets(slide, bullets, *, left, top, width, height,
                size: int = 18, color: str = "#222222") -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {b}"
        p.space_after = Pt(10)
        for run in p.runs:
            run.font.size = Pt(size)
            run.font.color.rgb = hex_color(color)


def add_rect(slide, *, left, top, width, height, fill: str) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(fill)
    shape.line.fill.background()


def set_notes(slide, notes: str) -> None:
    nf = slide.notes_slide.notes_text_frame
    nf.text = notes


# ----------------------------------------------------------------- renderers


def render_title(slide, palette, spec) -> None:
    set_background(slide, palette["background"])
    add_rect(slide, left=0, top=Inches(3.0), width=SLIDE_W, height=Inches(0.15),
             fill=palette["accent"])
    add_textbox(slide, spec.get("title", "Untitled"),
                left=Inches(0.75), top=Inches(2.0),
                width=Inches(11.8), height=Inches(1.2),
                size=44, bold=True, color=palette["primary"])
    if spec.get("subtitle"):
        add_textbox(slide, spec["subtitle"],
                    left=Inches(0.75), top=Inches(3.4),
                    width=Inches(11.8), height=Inches(0.8),
                    size=24, color=palette["text"])
    footer_bits = [spec.get("presenter"), spec.get("date")]
    footer = "  ·  ".join([s for s in footer_bits if s])
    if footer:
        add_textbox(slide, footer,
                    left=Inches(0.75), top=Inches(6.5),
                    width=Inches(11.8), height=Inches(0.5),
                    size=14, color=palette["text"])


def render_section(slide, palette, spec) -> None:
    set_background(slide, palette["primary"])
    add_textbox(slide, spec.get("title", ""),
                left=Inches(0.75), top=Inches(3.0),
                width=Inches(11.8), height=Inches(1.5),
                size=48, bold=True, color="#FFFFFF")
    add_rect(slide, left=Inches(0.75), top=Inches(4.6),
             width=Inches(2), height=Inches(0.1),
             fill=palette["accent"])


def _title_bar(slide, palette, title: str) -> None:
    add_rect(slide, left=0, top=0, width=SLIDE_W, height=Inches(1.0),
             fill=palette["primary"])
    add_textbox(slide, title,
                left=Inches(0.5), top=Inches(0.22),
                width=Inches(12.3), height=Inches(0.7),
                size=26, bold=True, color="#FFFFFF")


def render_content(slide, palette, spec) -> None:
    set_background(slide, palette["background"])
    _title_bar(slide, palette, spec.get("title", ""))
    bullets = spec.get("bullets") or []
    if bullets:
        add_bullets(slide, bullets,
                    left=Inches(0.75), top=Inches(1.5),
                    width=Inches(11.8), height=Inches(5.5),
                    size=20, color=palette["text"])
    if spec.get("body"):
        add_textbox(slide, spec["body"],
                    left=Inches(0.75), top=Inches(1.5),
                    width=Inches(11.8), height=Inches(5.5),
                    size=18, color=palette["text"])


def render_two_column(slide, palette, spec) -> None:
    set_background(slide, palette["background"])
    _title_bar(slide, palette, spec.get("title", ""))
    for col_key, col_left in (("left", Inches(0.5)), ("right", Inches(6.85))):
        col = spec.get(col_key, {})
        if col.get("title"):
            add_textbox(slide, col["title"],
                        left=col_left, top=Inches(1.3),
                        width=Inches(6), height=Inches(0.6),
                        size=20, bold=True, color=palette["primary"])
        if col.get("bullets"):
            add_bullets(slide, col["bullets"],
                        left=col_left, top=Inches(2.0),
                        width=Inches(6), height=Inches(5),
                        size=18, color=palette["text"])
        elif col.get("body"):
            add_textbox(slide, col["body"],
                        left=col_left, top=Inches(2.0),
                        width=Inches(6), height=Inches(5),
                        size=18, color=palette["text"])


def render_big_number(slide, palette, spec) -> None:
    set_background(slide, palette["background"])
    _title_bar(slide, palette, spec.get("title", ""))
    add_textbox(slide, spec.get("number", "—"),
                left=Inches(0.5), top=Inches(1.8),
                width=Inches(12.3), height=Inches(3.5),
                size=140, bold=True, color=palette["accent"],
                align=PP_ALIGN.CENTER)
    if spec.get("caption"):
        add_textbox(slide, spec["caption"],
                    left=Inches(0.5), top=Inches(5.6),
                    width=Inches(12.3), height=Inches(1),
                    size=22, color=palette["text"],
                    align=PP_ALIGN.CENTER)


def render_quote(slide, palette, spec) -> None:
    set_background(slide, palette["background"])
    add_rect(slide, left=Inches(1.4), top=Inches(2.0),
             width=Inches(0.12), height=Inches(3.5),
             fill=palette["accent"])
    add_textbox(slide, f"“{spec.get('quote', '')}”",
                left=Inches(1.9), top=Inches(2.0),
                width=Inches(10), height=Inches(3),
                size=30, color=palette["primary"])
    if spec.get("attribution"):
        add_textbox(slide, f"— {spec['attribution']}",
                    left=Inches(1.9), top=Inches(5.4),
                    width=Inches(10), height=Inches(0.7),
                    size=18, color=palette["text"])


def render_closing(slide, palette, spec) -> None:
    set_background(slide, palette["primary"])
    add_textbox(slide, spec.get("title", "Thank you"),
                left=Inches(0.5), top=Inches(2.5),
                width=Inches(12.3), height=Inches(1.5),
                size=54, bold=True, color="#FFFFFF",
                align=PP_ALIGN.CENTER)
    if spec.get("subtitle"):
        add_textbox(slide, spec["subtitle"],
                    left=Inches(0.5), top=Inches(4.3),
                    width=Inches(12.3), height=Inches(0.8),
                    size=22, color=palette["accent"],
                    align=PP_ALIGN.CENTER)


RENDERERS = {
    "title": render_title,
    "section": render_section,
    "content": render_content,
    "two_column": render_two_column,
    "big_number": render_big_number,
    "quote": render_quote,
    "closing": render_closing,
}


# ----------------------------------------------------------------- driver


def build(spec_path: Path, out_path: Path, script_dir: Path) -> None:
    spec = json.loads(spec_path.read_text())
    palette = resolve_palette(spec.get("palette"), script_dir)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    slides = spec.get("slides", [])
    if not slides:
        sys.stderr.write("error: spec has no slides\n")
        sys.exit(1)

    for i, slide_spec in enumerate(slides, start=1):
        layout = slide_spec.get("layout", "content")
        renderer = RENDERERS.get(layout)
        if renderer is None:
            sys.stderr.write(
                f"warning: unknown layout {layout!r} on slide {i}; falling back to content\n"
            )
            renderer = render_content
        slide = prs.slides.add_slide(blank_layout)
        renderer(slide, palette, slide_spec)
        if slide_spec.get("notes"):
            set_notes(slide, slide_spec["notes"])

    out_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out_path))
    print(f"wrote {out_path} ({len(slides)} slides)")


def main(argv) -> int:
    ap = argparse.ArgumentParser(description=__doc__, formatter_class=argparse.RawDescriptionHelpFormatter)
    ap.add_argument("spec", help="Path to deck spec JSON")
    ap.add_argument("-o", "--output", default=None,
                    help="Output .pptx path (default: alongside the spec)")
    args = ap.parse_args(argv)

    spec_path = Path(args.spec).expanduser().resolve()
    if not spec_path.exists():
        sys.stderr.write(f"error: spec file not found: {spec_path}\n")
        return 1

    out_path = (
        Path(args.output).expanduser().resolve()
        if args.output
        else spec_path.with_suffix(".pptx")
    )

    script_dir = Path(__file__).resolve().parent
    build(spec_path, out_path, script_dir)
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv[1:]))
